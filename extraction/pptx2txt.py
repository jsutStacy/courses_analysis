from db.DataModel import Lecture, db
import peewee
import os.path
from pptx import Presentation, exc
import pathos.multiprocessing as mp


class Pptx2Txt(object):
    def __init__(self, prefix, process_count=1):
        self.prefix = prefix
        self.pool = mp.ProcessingPool(process_count)

    def __convert(self, lecture):
        path = self.prefix+lecture.path
        if not os.path.exists(path):
            print "File not found: {0}".format(path)
            return
        print lecture.url

        ret_val = ''
        try:
            prs = Presentation(path)

            # text_runs will be populated with a list of strings,
            # one for each text run in presentation
            text_runs = []

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)

            ret_val = ' '.join(text_runs)

        except exc.PythonPptxError as e:
            print "Could not extract text {0}".format(e)

        lecture.content = ret_val
        return lecture

    def extract_text(self):
        lectures = Lecture.select().where(Lecture.content == '', Lecture.url % "*pptx")

        result_lectures = self.pool.map(self.__convert, lectures)

        for lecture in result_lectures:
            if lecture:
                try:
                    with db.transaction():
                        lecture.save()
                except peewee.OperationalError as e:
                    print e